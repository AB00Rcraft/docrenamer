"""Генераторы синтетического тестового корпуса (раздел 73 ТЗ).

Реальные конфиденциальные материалы не используются: все образцы создаются
программно и содержат вымышленные данные.
"""

from __future__ import annotations

import io
import json
import zipfile
from pathlib import Path

import pytest

DEJAVU = Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")

POSTANOVLENIE_TEXT = (
    "ПОСТАНОВЛЕНИЕ\n"
    "о возбуждении исполнительного производства\n"
    "№ 859189755/7728 от 27 июля 2026 года\n"
    "Алтуфьевский ОСП ГУФССП России по г. Москве\n"
    "Судебный пристав-исполнитель Сидорова А.А.\n"
    "Должник: Иванов Иван Иванович\n"
    "Взыскатель: ООО «Альфа», ИНН 7701234567\n"
    "Исполнительное производство № 652102/26/77028-ИП\n"
    "Сумма взыскания: 154 300,50 руб.\n"
)

DOGOVOR_TEXT = (
    "ДОГОВОР ЗАЙМА № 17\n"
    "город Москва, 18 августа 2026 года\n"
    "Общество с ограниченной ответственностью «Альфа», ИНН 7701234567, "
    "ОГРН 1157746000000, именуемое Займодавец, и Петров Сергей Андреевич, "
    "именуемый Заёмщик, заключили настоящий договор о нижеследующем.\n"
    "1. Займодавец передаёт Заёмщику 500 000 (пятьсот тысяч) рублей.\n"
    "2. Срок возврата — 18 августа 2027 года.\n"
)


def require(condition: bool, reason: str) -> None:
    """Пропустить тест, если необходимый инструмент недоступен."""
    if not condition:
        pytest.skip(reason)


# --- документы ------------------------------------------------------------


def make_pdf_with_text(path: Path, text: str = POSTANOVLENIE_TEXT) -> Path:
    """PDF с настоящим текстовым слоем на кириллице."""
    require(DEJAVU.is_file(), "Нет шрифта с кириллицей для генерации PDF")
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas

    pdfmetrics.registerFont(TTFont("DejaVu", str(DEJAVU)))
    pdf = canvas.Canvas(str(path), pagesize=A4)
    pdf.setTitle("Постановление")
    pdf.setAuthor("Алтуфьевский ОСП")
    pdf.setFont("DejaVu", 12)
    y = 800
    for line in text.splitlines():
        pdf.drawString(50, y, line)
        y -= 20
    pdf.showPage()
    pdf.save()
    return path


def make_pdf_scan(path: Path, text: str = "ПОСТАНОВЛЕНИЕ") -> Path:
    """PDF без текстового слоя: страница-изображение (имитация скана)."""
    require(DEJAVU.is_file(), "Нет шрифта с кириллицей для генерации скана")
    from PIL import Image, ImageDraw, ImageFont
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    image = Image.new("RGB", (1240, 1754), "white")
    draw = ImageDraw.Draw(image)
    font = ImageFont.truetype(str(DEJAVU), 40)
    for index, line in enumerate(text.splitlines()):
        draw.text((80, 100 + index * 60), line, fill="black", font=font)
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    buffer.seek(0)

    pdf = canvas.Canvas(str(path), pagesize=A4)
    pdf.drawImage(ImageReader(buffer), 0, 0, width=A4[0], height=A4[1])
    pdf.showPage()
    pdf.save()
    return path


def make_pdf_bad_text_layer(path: Path) -> Path:
    """PDF с непустым, но нечитаемым текстовым слоем.

    Так выглядит документ с испорченным отображением шрифта в Unicode: текст
    формально извлекается, но осмысленного русского в нём нет.
    """
    require(DEJAVU.is_file(), "Нет шрифта с кириллицей")
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas

    broken = POSTANOVLENIE_TEXT.encode("utf-8").decode("latin-1")
    pdfmetrics.registerFont(TTFont("DejaVu", str(DEJAVU)))
    pdf = canvas.Canvas(str(path), pagesize=A4)
    pdf.setFont("DejaVu", 11)
    y = 800
    for line in broken.splitlines():
        pdf.drawString(40, y, line[:120])
        y -= 18
    pdf.showPage()
    pdf.save()
    return path


def make_pdf_encrypted(path: Path) -> Path:
    """PDF, защищённый паролем."""
    require(DEJAVU.is_file(), "Нет шрифта с кириллицей")
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    pdf = canvas.Canvas(str(path), pagesize=A4, encrypt="секрет")
    pdf.drawString(50, 800, "secret")
    pdf.showPage()
    pdf.save()
    return path


def make_docx(path: Path, text: str = DOGOVOR_TEXT) -> Path:
    """DOCX с абзацами, таблицей и свойствами документа."""
    import docx

    document = docx.Document()
    document.core_properties.title = "Договор займа"
    document.core_properties.author = "Иванов И.И."
    for line in text.splitlines():
        document.add_paragraph(line)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Займодавец"
    table.cell(0, 1).text = "ООО «Альфа»"
    table.cell(1, 0).text = "Заёмщик"
    table.cell(1, 1).text = "Петров Сергей Андреевич"
    document.save(str(path))
    return path


def make_xlsx(path: Path) -> Path:
    """XLSX с русскими заголовками и данными."""
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Реестр"
    sheet.append(["Дата", "Контрагент", "Сумма", "Договор"])
    sheet.append(["18.08.2026", "ООО «Альфа»", 154300.5, "№ 17"])
    sheet.append(["19.08.2026", "ИП Петров С.А.", 25000, "№ 18"])
    workbook.properties.title = "Реестр платежей"
    workbook.properties.creator = "Бухгалтерия"
    workbook.save(str(path))
    return path


def make_pptx(path: Path) -> Path:
    """PPTX с заголовком слайда."""
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Отчёт по делу Иванова"
    slide.placeholders[1].text = "Исполнительное производство № 652102/26/77028-ИП"
    presentation.core_properties.title = "Отчёт"
    presentation.save(str(path))
    return path


def make_text(path: Path, text: str = DOGOVOR_TEXT, encoding: str = "utf-8") -> Path:
    """Текстовый файл в заданной кодировке."""
    path.write_bytes(text.encode(encoding))
    return path


def make_csv(path: Path, encoding: str = "cp1251") -> Path:
    """CSV с русскими заголовками и разделителем «;»."""
    content = "Дата;Контрагент;Сумма\n18.08.2026;ООО «Альфа»;154300,50\n"
    path.write_bytes(content.encode(encoding))
    return path


def make_html(path: Path, encoding: str = "windows-1251") -> Path:
    """HTML с объявленной кодировкой и скриптом, который нельзя исполнять."""
    html = (
        f'<html><head><meta charset="{encoding}">'
        "<title>Договор займа № 17</title></head><body>"
        "<h1>ДОГОВОР ЗАЙМА</h1>"
        "<script>window.location='http://example.invalid'</script>"
        "<p>Петров Сергей Андреевич, 18 августа 2026 года</p>"
        "</body></html>"
    )
    path.write_bytes(html.encode(encoding))
    return path


def make_xml_bomb(path: Path) -> Path:
    """XML с внешней сущностью — вход, который обязан быть отклонён."""
    payload = (
        '<?xml version="1.0"?>'
        '<!DOCTYPE doc [<!ENTITY xxe SYSTEM "file:///etc/passwd">]>'
        "<doc>&xxe;</doc>"
    )
    path.write_text(payload, encoding="utf-8")
    return path


def make_gpx(path: Path) -> Path:
    """GPX-трек с двумя точками."""
    payload = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<gpx version="1.1"><trk><name>Поездка Москва</name><trkseg>'
        '<trkpt lat="55.7558" lon="37.6173"><time>2026-08-03T10:00:00Z</time></trkpt>'
        '<trkpt lat="55.7658" lon="37.6273"><time>2026-08-03T11:00:00Z</time></trkpt>'
        "</trkseg></trk></gpx>"
    )
    path.write_text(payload, encoding="utf-8")
    return path


def make_json(path: Path) -> Path:
    """JSON с русскими ключами."""
    payload = {
        "наименование": "Договор займа",
        "дата": "2026-08-18",
        "номер": "17",
        "стороны": [{"фио": "Петров Сергей Андреевич", "роль": "Заёмщик"}],
    }
    path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    return path


def make_eml(path: Path) -> Path:
    """Письмо EML с закодированными русскими заголовками."""
    from email.message import EmailMessage

    message = EmailMessage()
    message["Subject"] = "Проект договора займа"
    message["From"] = "Иванов <ivanov@example.invalid>"
    message["To"] = "Петров <petrov@example.invalid>"
    message["Date"] = "Thu, 14 May 2026 10:30:00 +0300"
    message["Message-ID"] = "<msg-1@example.invalid>"
    message.set_content("Добрый день! Направляю проект договора займа № 17.")
    message.add_attachment(
        b"%PDF-1.4", maintype="application", subtype="pdf", filename="Договор.pdf"
    )
    path.write_bytes(message.as_bytes())
    return path


def make_rtf(path: Path) -> Path:
    """RTF с кириллицей в виде hex-escape-последовательностей CP1251."""
    text = "Договор займа"
    escaped = "".join(f"\\'{byte:02x}" for byte in text.encode("cp1251"))
    payload = r"{\rtf1\ansi\ansicpg1251{\fonttbl{\f0 Arial;}}\f0 " + escaped + "}"
    path.write_text(payload, encoding="ascii")
    return path


# --- изображения и медиа --------------------------------------------------


def make_jpeg(path: Path, size: tuple[int, int] = (1600, 1200)) -> Path:
    """JPEG нужного размера (без EXIF)."""
    from PIL import Image

    Image.new("RGB", size, (200, 210, 220)).save(str(path), format="JPEG")
    return path


def make_png_document(path: Path, text: str = "ДОГОВОР ЗАЙМА № 17") -> Path:
    """PNG со скан-подобным текстом для OCR."""
    require(DEJAVU.is_file(), "Нет шрифта с кириллицей")
    from PIL import Image, ImageDraw, ImageFont

    image = Image.new("RGB", (1240, 700), "white")
    draw = ImageDraw.Draw(image)
    font = ImageFont.truetype(str(DEJAVU), 44)
    for index, line in enumerate(text.splitlines()):
        draw.text((60, 60 + index * 70), line, fill="black", font=font)
    image.save(str(path))
    return path


def make_wav(path: Path, seconds: float = 1.5) -> Path:
    """WAV-файл заданной длительности."""
    import struct
    import wave

    with wave.open(str(path), "wb") as handle:
        handle.setnchannels(1)
        handle.setsampwidth(2)
        handle.setframerate(8000)
        frames = int(8000 * seconds)
        handle.writeframes(b"".join(struct.pack("<h", 0) for _ in range(frames)))
    return path


# --- архивы ---------------------------------------------------------------


def make_mp4(
    path: Path,
    *,
    created: str = "2026-08-12T17:48:22",
    duration_seconds: float = 102.0,
    location: str = "+55.7558+37.6173/",
) -> Path:
    """Минимальный контейнер MP4 с корректными боксами ftyp/moov.

    Реального видеопотока внутри нет: проверяется разбор метаданных, а не
    декодирование.
    """
    import struct
    from datetime import UTC, datetime

    def box(kind: bytes, payload: bytes) -> bytes:
        return struct.pack(">I", len(payload) + 8) + kind + payload

    epoch = datetime(1904, 1, 1, tzinfo=UTC)
    moment = datetime.fromisoformat(created).replace(tzinfo=UTC)
    seconds = int((moment - epoch).total_seconds())
    timescale = 1000
    duration = int(duration_seconds * timescale)

    mvhd = box(
        b"mvhd",
        struct.pack(">BBBB", 0, 0, 0, 0)
        + struct.pack(">IIII", seconds, seconds, timescale, duration)
        + b"\x00" * 80,
    )
    hdlr = box(b"hdlr", b"\x00" * 8 + b"vide" + b"\x00" * 12)
    mdia = box(b"mdia", hdlr)
    trak = box(b"trak", mdia)
    xyz = box(b"\xa9xyz", struct.pack(">HH", len(location), 0x15C7) + location.encode())
    udta = box(b"udta", xyz)
    moov = box(b"moov", mvhd + trak + udta)
    ftyp = box(b"ftyp", b"isom" + struct.pack(">I", 512) + b"isomiso2mp41")

    path.write_bytes(ftyp + moov)
    return path


def make_jpeg_with_exif(
    path: Path,
    *,
    make: str = "Apple",
    model: str = "iPhone 16 Pro",
    taken: str = "2026:08:03 18:42:17",
    latitude: tuple[int, int, float] = (55, 45, 20.88),
    longitude: tuple[int, int, float] = (37, 37, 2.28),
) -> Path:
    """JPEG с настоящим EXIF: датой съёмки, устройством и координатами."""
    from PIL import Image
    from PIL.ExifTags import GPS, Base
    from PIL.TiffImagePlugin import IFDRational

    image = Image.new("RGB", (4032, 3024), (180, 190, 200))
    exif = image.getexif()
    exif[Base.Make.value] = make
    exif[Base.Model.value] = model
    exif[Base.DateTimeOriginal.value] = taken
    gps = exif.get_ifd(0x8825)
    gps[GPS.GPSLatitudeRef.value] = "N"
    gps[GPS.GPSLatitude.value] = tuple(IFDRational(round(v * 100), 100) for v in latitude)
    gps[GPS.GPSLongitudeRef.value] = "E"
    gps[GPS.GPSLongitude.value] = tuple(IFDRational(round(v * 100), 100) for v in longitude)
    image.save(str(path), exif=exif)
    return path


def make_zip(path: Path, names: list[str] | None = None) -> Path:
    """ZIP со списком файлов (архив не распаковывается при анализе)."""
    names = names or [
        "договоры/Договор-Альфа-1.pdf",
        "договоры/Договор-Альфа-2.pdf",
        "договоры/Договор-Альфа-3.pdf",
    ]
    with zipfile.ZipFile(path, "w") as archive:
        for name in names:
            archive.writestr(name, "содержимое")
    return path


def make_zip_bomb(path: Path) -> Path:
    """ZIP с экстремальным коэффициентом сжатия."""
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("bomb.txt", "0" * 5_000_000)
    return path


def make_corrupted(path: Path, kind: str = "pdf") -> Path:
    """Файл с корректной сигнатурой и мусорным содержимым."""
    signatures = {
        "pdf": b"%PDF-1.7\n",
        "zip": b"PK\x03\x04",
        "jpg": b"\xff\xd8\xff\xe0",
    }
    path.write_bytes(signatures.get(kind, b"") + bytes(range(256)) * 4)
    return path

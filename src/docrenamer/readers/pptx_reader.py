"""PPTX (раздел 20 ТЗ). Извлекаются слайды, заголовки, заметки и свойства."""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING

from docrenamer.readers.base import finalize_text, safe_metadata
from docrenamer.readers.docx_reader import _core_properties
from docrenamer.types import ReadResult, Status

if TYPE_CHECKING:  # pragma: no cover
    from docrenamer.analysis import ReaderContext

MAX_SLIDES = 300


def read_pptx(path: Path, context: ReaderContext) -> ReadResult:
    """Прочитать презентацию."""
    result = ReadResult()
    limits = context.limits
    try:
        from pptx import Presentation
    except ImportError:  # pragma: no cover
        result.add_status(Status.UNSUPPORTED_FORMAT)
        return result

    try:
        presentation = Presentation(str(path))
    except Exception as exc:  # недоверенный вход
        result.add_status(Status.READ_ERROR)
        result.decoding_warnings.append(f"Презентация не открыта: {exc}")
        return result

    parts: list[str] = []
    titles: list[str] = []
    notes: list[str] = []
    slides = list(presentation.slides)[:MAX_SLIDES]
    for index, slide in enumerate(slides, start=1):
        parts.append(f"[Слайд {index}]")
        try:
            title_shape = slide.shapes.title
        except (AttributeError, ValueError):
            title_shape = None
        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            title = title_shape.text_frame.text.strip()
            if title:
                titles.append(title)
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)
        try:
            if slide.has_notes_slide:
                note = slide.notes_slide.notes_text_frame.text.strip()
                if note:
                    notes.append(note)
        except (AttributeError, ValueError, KeyError):
            continue

    result.metadata.update(_core_properties(presentation))
    result.metadata.update(
        safe_metadata({"slide_count": len(slides), "slide_titles": titles[:20]})
    )
    result.source_encoding = "ooxml/utf-8"
    result.encoding_confidence = 1.0
    if notes:
        parts.extend(notes[:20])
    return finalize_text(result, "\n".join(parts), limits)
